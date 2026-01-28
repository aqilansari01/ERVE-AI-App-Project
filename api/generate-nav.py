from http.server import BaseHTTPRequestHandler
import json
import base64
import re
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import traceback

class handler(BaseHTTPRequestHandler):
    def set_cell_text(self, cell, text, font_name='Arial', font_size=8, bold=False):
        """Helper to set cell text with proper formatting"""
        cell.text = str(text)
        # Set font for all paragraphs and runs in the cell
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = bold

    def do_POST(self):
        try:
            content_length = int(self.headers['Content-Length'])
            post_data = self.rfile.read(content_length)
            data = json.loads(post_data.decode('utf-8'))

            template_base64 = data.get('templateFile')
            rag_status = data.get('ragStatus', {})
            quarterly_financials = data.get('quarterlyFinancials', {})
            company_update = data.get('companyUpdate', '')
            investment_summary = data.get('investmentSummary', {})

            # Debug logging
            print("="*60)
            print("DATA RECEIVED FROM FRONTEND:")
            print("="*60)
            print("Investment Summary:", investment_summary)
            print("RAG Status:", rag_status)
            print("Quarterly Financials:", quarterly_financials)
            print("Company Update length:", len(company_update) if company_update else 0)
            print("="*60)

            if not template_base64:
                self.send_error(400, 'No template file provided')
                return

            template_bytes = base64.b64decode(template_base64)
            template_stream = BytesIO(template_bytes)
            prs = Presentation(template_stream)

            if len(prs.slides) == 0:
                self.send_error(400, 'Template has no slides')
                return

            slide = prs.slides[0]
            updates = []

            # 1. Update Investment Summary table with manual inputs
            inv_updated = self.update_investment_summary_table(slide, investment_summary)
            if inv_updated:
                updates.append(f"Updated investment summary table")

            # 2. Update RAG status indicators
            rag_updated = self.update_rag_status(slide, rag_status)
            if rag_updated:
                updates.append(f"Updated {rag_updated} RAG indicators")

            # 3. Update quarterly financials table
            fin_updated = self.update_quarterly_financials(slide, quarterly_financials)
            if fin_updated:
                updates.append(f"Updated quarterly financials table")

            # 4. Update company update text
            update_updated = self.update_company_update(slide, company_update)
            if update_updated:
                updates.append(f"Updated company update section")

            # Note: Exit cases and valuation waterfall are left unchanged as requested

            output_stream = BytesIO()
            prs.save(output_stream)
            output_stream.seek(0)

            output_base64 = base64.b64encode(output_stream.read()).decode('utf-8')

            self.send_response(200)
            self.send_header('Content-type', 'application/json')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.end_headers()

            response = {
                'success': True,
                'file': output_base64,
                'updates': updates,
                'message': f"Successfully updated: {', '.join(updates)}" if updates else "No updates made"
            }
            self.wfile.write(json.dumps(response).encode('utf-8'))

        except Exception as e:
            error_details = traceback.format_exc()
            self.send_response(500)
            self.send_header('Content-type', 'application/json')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.end_headers()
            error_response = {
                'success': False,
                'error': str(e),
                'details': error_details
            }
            self.wfile.write(json.dumps(error_response).encode('utf-8'))

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def update_investment_summary_table(self, slide, investment_summary):
        """Update the Investment Summary table with manual input data"""
        if not investment_summary:
            return False

        try:
            # Find the Investment Summary table - it's the first table after "Investment summary" text
            found_investment_header = False

            for shape in slide.shapes:
                # First find the "Investment summary" text
                if shape.has_text_frame and 'investment summary' in shape.text.lower():
                    found_investment_header = True
                    continue

                # If we found the header, the next table is our target
                if found_investment_header and shape.has_table:
                    table = shape.table

                    # Build governance string based on what's available
                    board_member = investment_summary.get('boardMember', '')
                    board_observer = investment_summary.get('boardObserver', '')
                    governance_parts = []
                    if board_member:
                        governance_parts.append(f"Board Member: {board_member}")
                    if board_observer:
                        governance_parts.append(f"Observer: {board_observer}")
                    governance_value = ', '.join(governance_parts) if governance_parts else ''

                    # Map field labels to values
                    field_map = {
                        'erve investment': investment_summary.get('erveInvestmentEUR', ''),
                        'break-down by round': investment_summary.get('investmentRound', ''),
                        'breakdown by round': investment_summary.get('investmentRound', ''),
                        'total raised': investment_summary.get('totalRaised', ''),
                        'erve %': investment_summary.get('erveOwnership', ''),
                        'security': investment_summary.get('securityType', ''),
                        'other shareholders': investment_summary.get('otherShareholders', ''),
                        'governance': governance_value,
                        # Note: 'cash' field is not in the form, leave blank in template
                        'monthly burn': investment_summary.get('monthlyBurn', ''),
                        'fume': investment_summary.get('fume', ''),
                        'last pre-/post-money valuation': f"€{investment_summary.get('preMoneyValuation', '')}m / €{investment_summary.get('postMoneyValuation', '')}m" if investment_summary.get('preMoneyValuation') else '',
                    }

                    # Add NAV values with flexible quarter matching
                    # These will match any cell containing 'nav' (case insensitive)
                    current_nav_eur = investment_summary.get('currentQuarterNAV', '')
                    current_nav_usd = investment_summary.get('currentQuarterNAVUSD', '')
                    prior_nav_eur = investment_summary.get('priorQuarterNAV', '')
                    prior_nav_usd = investment_summary.get('priorQuarterNAVUSD', '')

                    # Track NAV cells to update them in order
                    nav_cells = []

                    print("\n=== Processing Investment Summary Table ===")
                    print(f"Field values to update:")
                    for field_label, field_value in field_map.items():
                        print(f"  {field_label}: {repr(field_value)}")

                    # Iterate through all cells to find and update matching fields
                    for row in table.rows:
                        for i, cell in enumerate(row.cells):
                            cell_text = cell.text.strip().lower()

                            # Special handling for NAV fields (flexible quarter matching)
                            if 'nav' in cell_text and i + 1 < len(row.cells):
                                nav_cells.append((cell_text, row.cells[i + 1]))
                                continue

                            # Check if this cell contains a field label
                            for field_label, field_value in field_map.items():
                                if field_label in cell_text:
                                    print(f"Found label '{field_label}' in cell: '{cell_text}'")
                                    # Update the cell to the right (next cell in row)
                                    if i + 1 < len(row.cells):
                                        if field_value or field_value == 0:  # Allow 0 but not empty string
                                            # For ERVE investment, format as €Xm
                                            if field_label == 'erve investment':
                                                self.set_cell_text(row.cells[i + 1], f"€{field_value}m")
                                                print(f"  ✓ Updated ERVE investment: €{field_value}m")
                                            # For ERVE %, add %
                                            elif field_label == 'erve %':
                                                self.set_cell_text(row.cells[i + 1], f"{field_value}%")
                                                print(f"  ✓ Updated ERVE %: {field_value}%")
                                            # For monthly burn, format as €Xm
                                            elif field_label == 'monthly burn':
                                                self.set_cell_text(row.cells[i + 1], f"€{field_value}m")
                                                print(f"  ✓ Updated monthly burn: €{field_value}m")
                                            # For FUME, add "months"
                                            elif field_label == 'fume':
                                                self.set_cell_text(row.cells[i + 1], f"{field_value} months")
                                                print(f"  ✓ Updated FUME: {field_value} months")
                                            else:
                                                self.set_cell_text(row.cells[i + 1], str(field_value))
                                                print(f"  ✓ Updated {field_label}: {field_value}")
                                        else:
                                            print(f"  ✗ Skipped {field_label}: value is empty ({repr(field_value)})")
                                    break

                    # Update NAV cells with smarter matching
                    # Sort NAV cells to handle them in a consistent order (top to bottom)
                    for label_text, value_cell in nav_cells:
                        # Determine if EUR or USD based on label
                        is_eur = '(eur)' in label_text or '€' in label_text or 'eur' in label_text
                        is_usd = '(usd)' in label_text or '$' in label_text or 'usd' in label_text

                        # If neither EUR nor USD is explicitly mentioned, assume first occurrence is EUR
                        if not is_eur and not is_usd:
                            nav_index = nav_cells.index((label_text, value_cell))
                            is_eur = nav_index % 2 == 0  # Even indices are EUR
                            is_usd = nav_index % 2 == 1  # Odd indices are USD

                        # Try to determine if current or prior quarter
                        # Check for quarter indicators in the label
                        nav_index = nav_cells.index((label_text, value_cell))

                        # First try specific quarter matching (Q4, Q3, Q2, Q1, etc.)
                        # Extract quarter number if present
                        quarter_match = re.search(r'q(\d)', label_text)

                        # Determine if this is current or prior based on position and quarter
                        # Assume NAV cells appear in order: most recent first
                        is_current = nav_index < len(nav_cells) / 2

                        # Update the cell based on currency and position
                        if is_current:
                            if is_eur and current_nav_eur:
                                self.set_cell_text(value_cell, current_nav_eur)
                                print(f"Updated current NAV EUR: {current_nav_eur}")
                            elif is_usd and current_nav_usd:
                                self.set_cell_text(value_cell, current_nav_usd)
                                print(f"Updated current NAV USD: {current_nav_usd}")
                        else:
                            if is_eur and prior_nav_eur:
                                self.set_cell_text(value_cell, prior_nav_eur)
                                print(f"Updated prior NAV EUR: {prior_nav_eur}")
                            elif is_usd and prior_nav_usd:
                                self.set_cell_text(value_cell, prior_nav_usd)
                                print(f"Updated prior NAV USD: {prior_nav_usd}")

                    return True

            return False

        except Exception as e:
            print(f"Error updating investment summary table: {e}")
            traceback.print_exc()
            return False

    def update_rag_status(self, slide, rag_status):
        """Update RAG status indicators by finding colored rectangles"""
        try:
            rag_colors = {
                'Green': RGBColor(77, 191, 175),
                'Amber': RGBColor(255, 192, 0),
                'Red': RGBColor(192, 0, 0),
            }

            updated_count = 0
            rectangles = []

            for shape in slide.shapes:
                if hasattr(shape, 'shape_type'):
                    if shape.shape_type == 1:  # Auto shape
                        width = shape.width
                        height = shape.height
                        # RAG indicators are small rectangles
                        if height < Inches(0.5) and width < Inches(1.5):
                            rectangles.append(shape)

            rectangles.sort(key=lambda s: (s.top, s.left))

            rag_values = [
                rag_status.get('financialsStatus', 'Green'),
                rag_status.get('cashStatus', 'Green'),
                rag_status.get('marketStatus', 'Green'),
                rag_status.get('teamStatus', 'Green'),
                rag_status.get('governanceStatus', 'Green'),
                rag_status.get('overallStatus', 'Green')
            ]

            for i, rect in enumerate(rectangles[:6]):
                if i < len(rag_values):
                    status = rag_values[i]
                    color = rag_colors.get(status, RGBColor(128, 128, 128))
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = color
                    updated_count += 1

            return updated_count

        except Exception as e:
            print(f"Error updating RAG status: {e}")
            traceback.print_exc()
            return 0

    def update_quarterly_financials(self, slide, quarterly_financials):
        """Update the quarterly financials table"""
        if not quarterly_financials:
            print("No quarterly financials data provided")
            return False

        print(f"Quarterly financials data received: {quarterly_financials}")

        try:
            # Find the table with "Quarterly Actuals" or metric names
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table

                    # Check if any cell contains "YF 31st Dec" or "Quarterly Actuals"
                    is_financials_table = False
                    for row in table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip().lower()
                            if 'yf 31st dec' in cell_text or 'quarterly actual' in cell_text:
                                print(f"Found financials table with text: {cell_text}")
                                is_financials_table = True
                                break
                        if is_financials_table:
                            break

                    if not is_financials_table:
                        # Also check for metric names in first column
                        for row in table.rows:
                            first_cell = row.cells[0].text.strip().upper()
                            if first_cell in ['ARR', 'REVENUE', 'GM', 'EBITDA', 'FTES']:
                                print(f"Found financials table with metric: {first_cell}")
                                is_financials_table = True
                                break

                    if not is_financials_table:
                        continue

                    # This is the financials table - find period columns
                    header_row = table.rows[0]
                    period_columns = {}

                    print("Header row cells:", [cell.text.strip() for cell in header_row.cells])
                    print("Available periods in data:", list(quarterly_financials.keys()))

                    for col_idx, cell in enumerate(header_row.cells):
                        cell_text = cell.text.strip()
                        for period in quarterly_financials.keys():
                            if period in cell_text:
                                period_columns[col_idx] = period
                                print(f"Matched column {col_idx} ({cell_text}) to period {period}")
                                break

                    if not period_columns:
                        print("WARNING: No period columns matched!")
                        return False

                    # Update data rows
                    metrics_map = {
                        'ARR': 'ARR',
                        'REVENUE': 'Revenue',
                        'GM': 'GM',
                        'EBITDA': 'EBITDA',
                        'FTES': 'FTEs'
                    }

                    updates_made = 0
                    for row in table.rows[1:]:
                        first_cell = row.cells[0].text.strip().upper()

                        if first_cell in metrics_map:
                            metric_name = metrics_map[first_cell]
                            print(f"Processing metric: {metric_name}")

                            for col_idx, period in period_columns.items():
                                if col_idx < len(row.cells):
                                    value = quarterly_financials[period].get(metric_name)
                                    if value is not None:
                                        self.set_cell_text(row.cells[col_idx], str(value))
                                        print(f"  Updated {metric_name} for {period}: {value}")
                                        updates_made += 1

                    print(f"Total cell updates in quarterly financials: {updates_made}")
                    return updates_made > 0

            print("No financials table found")
            return False

        except Exception as e:
            print(f"Error updating quarterly financials: {e}")
            traceback.print_exc()
            return False

    def update_company_update(self, slide, company_update):
        """Update the company update text box in the gray box"""
        if not company_update:
            print("No company update data provided")
            return False

        print(f"Company update data received (length: {len(company_update)} chars)")

        try:
            # Find text boxes in gray boxes (usually have gray fill)
            shapes_checked = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shapes_checked += 1
                    text_lower = shape.text.lower()

                    # Look for "company update" text
                    if 'company update' in text_lower:
                        print(f"Found 'company update' text box (shape type: {shape.shape_type})")
                        text_frame = shape.text_frame

                        # Clear existing text - safer approach
                        # Clear text from all existing paragraphs
                        for paragraph in text_frame.paragraphs:
                            paragraph.clear()

                        # Use first paragraph for header
                        if len(text_frame.paragraphs) > 0:
                            p1 = text_frame.paragraphs[0]
                        else:
                            p1 = text_frame.add_paragraph()

                        p1.text = "Company update"
                        # Set font on runs, not paragraph
                        for run in p1.runs:
                            run.font.name = 'Arial'
                            run.font.bold = True
                            run.font.size = Pt(8)

                        # Add content in second paragraph
                        clean_update = company_update.replace('## Company Update', '').replace('## Company update', '').strip()
                        p2 = text_frame.add_paragraph()
                        p2.text = clean_update
                        p2.space_before = Pt(6)
                        # Set font on runs
                        for run in p2.runs:
                            run.font.name = 'Arial'
                            run.font.size = Pt(7)

                        print(f"Successfully updated company update section")
                        return True

            print(f"Checked {shapes_checked} shapes with text frames, but 'company update' not found")
            return False

        except Exception as e:
            print(f"Error updating company update: {e}")
            traceback.print_exc()
            return False
